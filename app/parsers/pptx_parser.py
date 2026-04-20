"""Extract text and images from PPTX slides."""

import subprocess
import tempfile
from pathlib import Path

from pdf2image import convert_from_path
from pptx import Presentation


def extract_slides(pptx_path: str) -> list[str]:
    """Return a list of strings, one per slide, from the given PPTX file."""
    prs = Presentation(pptx_path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        slides.append("\n".join(texts))
    return slides


def extract_slide_images(pptx_path: str) -> list:
    """Return a list of PIL Images, one per slide.

    Tries LibreOffice first. If not installed, falls back to rendering each slide
    as a plain white image with the extracted text drawn on it.
    """
    libreoffice = subprocess.run(["which", "libreoffice"], capture_output=True).returncode == 0
    if libreoffice:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
                check=True,
                capture_output=True,
            )
            pdf_path = Path(tmpdir) / (Path(pptx_path).stem + ".pdf")
            return convert_from_path(str(pdf_path), dpi=150)

    # Fallback: render text-only images using Pillow
    from PIL import Image, ImageDraw, ImageFont
    slide_texts = extract_slides(pptx_path)
    images = []
    for text in slide_texts:
        img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
        except Exception:
            font = ImageFont.load_default()
        draw.multiline_text((60, 60), text, fill=(20, 40, 80), font=font, spacing=10)
        images.append(img)
    return images
