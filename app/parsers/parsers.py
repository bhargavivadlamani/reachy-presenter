def parse(file_path: str, parser: str = "pdfplumber") -> list[str]:
    if parser == "pdfplumber":
        return _parse_pdfplumber(file_path)
    elif parser == "python-pptx":
        return _parse_python_pptx(file_path)
    elif parser == "docling":
        return _parse_docling(file_path)
    raise ValueError(f"Unknown parser: {parser}")

def _parse_docling(path: str) -> list[str]:
    from docling.document_converter import DocumentConverter
    from docling.chunking import HybridChunker
    from docling.datamodel.pipeline_options import PdfPipelineOptions

    pipeline_options = PdfPipelineOptions(do_ocr=False)
    converter = DocumentConverter(
        format_options={"pdf": {"pipeline_options": pipeline_options}}
    )
    result = converter.convert(path)

    if result.status.name == "FAILURE":
        raise RuntimeError(f"Docling failed to convert {path}: {result.status}")

    chunker = HybridChunker()
    pages: dict[int, list[str]] = {}
    for chunk in chunker.chunk(result.document):
        text = chunker.contextualize(chunk)
        page_nos = {prov.page_no for item in chunk.meta.doc_items for prov in item.prov}
        page = min(page_nos) if page_nos else 0
        pages.setdefault(page, []).append(text)

    if not pages:
        return []

    max_page = max(pages.keys())
    return ["\n\n".join(pages.get(i, [])) for i in range(1, max_page + 1)]

def _parse_pdfplumber(path: str) -> list[str]:
    import pdfplumber
    with pdfplumber.open(path) as pdf:
        return [page.extract_text() or "" for page in pdf.pages]


def _parse_python_pptx(path: str) -> list[str]:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        text = "\n".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        slides.append(text)
    return slides
